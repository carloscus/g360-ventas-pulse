# -*- coding: utf-8 -*-
"""Manual PPTX de Ventas Pulse CIPSA — contenido derivado del README.md.
Widescreen 16:9 (13.333 x 7.5 in), branding CIPSA (verde 008F5D)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

W, H = 13.333, 7.5
VERDE = RGBColor(0x00, 0x8F, 0x5D)
VERDE_OSCURO = RGBColor(0x0B, 0x6B, 0x47)
TEXTO = RGBColor(0x1F, 0x29, 0x37)
GRIS = RGBColor(0x64, 0x74, 0x8B)
FONDO = RGBColor(0xFF, 0xFF, 0xFF)
SUAVE = RGBColor(0xE3, 0xF0, 0xEB)
AMBAR = RGBColor(0xB4, 0x53, 0x09)

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]
TOTAL = 10
CW = W - 1.2  # ancho util de contenido (margenes 0.6)

def nueva():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, Inches(W), Inches(H))
    bg.fill.solid(); bg.fill.fore_color.rgb = FONDO
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def caja(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf

def p(tf, text, size=14, bold=False, color=TEXTO, first=False, space=6, align=PP_ALIGN.LEFT):
    par = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    par.space_after = Pt(space)
    par.alignment = align
    r = par.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = 'Arial'
    return par

def seccion(s, titulo, num):
    tf = caja(s, 0.6, 0.3, 11.5, 0.5)
    p(tf, titulo, 21, True, VERDE_OSCURO, first=True)
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(0.85), Inches(CW), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = VERDE; ln.line.fill.background(); ln.shadow.inherit = False
    ft = caja(s, 10.2, 7.1, 2.55, 0.3)
    par = ft.paragraphs[0]; par.alignment = PP_ALIGN.RIGHT
    r = par.add_run(); r.text = f'powered by G360 · {num}/{TOTAL}'
    r.font.size = Pt(8); r.font.color.rgb = GRIS; r.font.name = 'Arial'

def tarjeta(s, x, y, w, h, titulo, cuerpo, tcolor=VERDE_OSCURO):
    card = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = SUAVE
    card.line.color.rgb = VERDE; card.line.width = Pt(1)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.1)
    p(tf, titulo, 12, True, tcolor, first=True, space=3)
    for linea in cuerpo:
        p(tf, linea, 9.5, False, TEXTO, space=2)

SHOTS = 'assets/screenshots'

def marco(s, x, y, w, h, archivo, etiqueta):
    """Captura de pantalla con marco telefono. Si existe la imagen se coloca;
    si no, dibuja el placeholder guia con la instruccion de captura."""
    from pathlib import Path
    ruta = Path(SHOTS) / archivo
    phone = s.shapes.add_shape(5, Inches(x - 0.07), Inches(y - 0.07), Inches(w + 0.14), Inches(h + 0.14))
    phone.fill.solid(); phone.fill.fore_color.rgb = TEXTO
    phone.line.color.rgb = VERDE_OSCURO; phone.line.width = Pt(1.75); phone.shadow.inherit = False
    if ruta.exists():
        # Encaje proporcional (contain) centrado: nunca distorsiona la captura
        from PIL import Image
        iw, ih = Image.open(ruta).size
        ir = iw / ih
        fr = w / h
        if ir > fr:
            dw, dh = w, w / ir
        else:
            dh, dw = h, h * ir
        dx = x + (w - dw) / 2
        dy = y + (h - dh) / 2
        s.shapes.add_picture(str(ruta), Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    else:
        ph = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xF9)
        ph.line.color.rgb = VERDE; ph.line.width = Pt(1)
        ph.line.dash_style = 4  # MSO_LINE_DASH_STYLE.DASH
        ph.shadow.inherit = False
        tf = ph.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
        p(tf, 'CAPTURA PENDIENTE', 8.5, True, GRIS, first=True, space=4, align=PP_ALIGN.CENTER)
        p(tf, etiqueta, 9, False, TEXTO, align=PP_ALIGN.CENTER)
        p(tf, f'Guardar como: {ruta.as_posix()}', 8, False, GRIS, align=PP_ALIGN.CENTER)

# ===== 1. PORTADA =====
s = nueva()
banda = s.shapes.add_shape(1, 0, Inches(3.15), Inches(W), Pt(4))
banda.fill.solid(); banda.fill.fore_color.rgb = VERDE; banda.line.fill.background(); banda.shadow.inherit = False
tf = caja(s, 1, 1.5, 11.3, 1.3)
p(tf, 'Ventas Pulse CIPSA', 44, True, VERDE_OSCURO, first=True, align=PP_ALIGN.CENTER)
tf = caja(s, 1.5, 3.5, 10.3, 1.2)
p(tf, 'Cockpit de campo para vendedores: PWA estática (SvelteKit + GitHub Pages)', 16, False, TEXTO, first=True, align=PP_ALIGN.CENTER, space=3)
p(tf, 'Dashboard diario · Radar de recompra · Montos netos · Ficha comercial · Búsqueda global', 12, False, GRIS, align=PP_ALIGN.CENTER)
tf = caja(s, 1, 5.4, 11.3, 0.9)
p(tf, 'Manual de uso · v2.0.0', 13, True, GRIS, first=True, align=PP_ALIGN.CENTER, space=2)
p(tf, 'Instalable en Android/iOS · Funciona offline · powered by G360', 10.5, False, GRIS, align=PP_ALIGN.CENTER)

# ===== 2. INTRODUCCION =====
s = nueva(); seccion(s, 'INTRODUCCIÓN', 2)
tf = caja(s, 0.6, 1.05, 12.1, 1.0)
p(tf, '¿Qué es Ventas Pulse?', 15, True, TEXTO, first=True, space=3)
p(tf, 'App web instalable (PWA) para el vendedor de campo de CIPSA: agenda del día, oportunidades de recompra, montos netos jerárquicos y ficha comercial por cliente, con stock en tiempo real y búsqueda global. Uso one-handed en móvil, datos del ERP vía Supabase, cache offline de 4 capas.', 11.5)
tarjeta(s, 0.6, 2.55, 5.95, 2.0, 'Para el vendedor', [
    'Prioridades del día al abrir la app (Hoy)',
    'Radar: qué cliente recomprar hoy y con qué argumentos',
    'Comparativa de precios entre sus clientes',
    'Ruta del día marcable con valor estimado'])
tarjeta(s, 6.75, 2.55, 5.95, 2.0, 'Bajo el capó', [
    'SvelteKit SPA estática en GitHub Pages · Supabase PostgREST (anon, solo lectura)',
    'Stock API (Render, refresh ~15 min) + catálogo estático',
    'Cache: memoria → sessionStorage → localStorage → red'])
tf = caja(s, 0.6, 4.85, 12.1, 2.2)
p(tf, 'Seguridad de la información', 13, True, AMBAR, first=True, space=3)
p(tf, 'La ficha solo abre clientes de tu cartera; el login se bloquea 3 minutos tras 2 intentos fallidos. Son frenos de cliente: la garantía real de confidencialidad por vendedor llega con Auth + RLS (Fase 4).', 11)

# ===== 3. MODULOS =====
s = nueva(); seccion(s, 'MÓDULOS', 3)
filas = [
    ('Hoy', '/dashboard', 'Agenda del día: prioritarios, próximos a recompra, alertas y "Cómo voy"'),
    ('Radar', '/radar', 'Oportunidades VENCIDAS priorizadas por valor estimado + ruta del día'),
    ('Netos', '/netos', 'Montos netos jerárquicos (cliente → línea → SKU), 3 periodos A/B/C'),
    ('Clientes', '/clientes', 'Directorio del vendedor (180 días, offline-first)'),
    ('Ficha', '/ficha', 'Precios por año con anomalías, comparativa Pareto y export XLSX'),
]
y = 1.1
for i, (mod, ruta, desc) in enumerate(filas):
    card = s.shapes.add_shape(5, Inches(0.6), Inches(y), Inches(8.1), Inches(0.98))
    card.fill.solid(); card.fill.fore_color.rgb = SUAVE if i % 2 == 0 else FONDO
    card.line.color.rgb = VERDE; card.line.width = Pt(0.75); card.shadow.inherit = False
    t1 = caja(s, 0.8, y + 0.1, 1.55, 0.4); p(t1, mod, 12.5, True, VERDE_OSCURO, first=True)
    t2 = caja(s, 2.35, y + 0.13, 1.8, 0.4); p(t2, ruta, 9.5, False, GRIS, first=True)
    t3 = caja(s, 4.2, y + 0.09, 4.4, 0.85); p(t3, desc, 9.5, False, TEXTO, first=True)
    y += 1.12
tf = caja(s, 0.6, y + 0.02, 8.1, 0.6)
p(tf, 'Navegación inferior fija: Hoy · Radar · Netos · Clientes. La ficha se abre desde Clientes o el buscador (lupa en el header).', 9.5, False, GRIS, first=True)
marco(s, 9.15, 1.1, 1.85, 4.0, 'radar.png',
      'Radar con 2-3 oportunidades y una marcada "En ruta"')
marco(s, 11.2, 1.1, 1.85, 4.0, 'netos.png',
      'Netos con 1 cliente expandido (lineas y variacion)')

# ===== 4. FLUJO DE TRABAJO =====
s = nueva(); seccion(s, 'FLUJO DE TRABAJO DIARIO', 4)
pasos = [
    ('1 · Ingresar', 'Login con tu código de vendedor + un cliente de tu cartera. La app valida el par y restaura la sesión.'),
    ('2 · Revisar Hoy', 'El dashboard lista prioritarios, próximos a vencer y tu avance "Cómo voy".'),
    ('3 · Planificar ruta', 'En Radar, marca con "+ Agregar" los clientes a visitar. La ruta queda en el dispositivo.'),
    ('4 · Preparar visita', 'Abre la ficha: KPIs, evolución mensual, precios por año y anomalías.'),
    ('5 · Negociar', '"Comparar precios entre mis clientes" por SKU: quién paga más/menos y tu posición.'),
    ('6 · Cerrar el ciclo', 'Marca la visita hecha en ruta. Todo queda cacheado para zonas sin señal.'),
]
y = 1.1
for i, (t, d) in enumerate(pasos):
    par_num = s.shapes.add_shape(9, Inches(0.6), Inches(y), Inches(0.44), Inches(0.44))
    par_num.fill.solid(); par_num.fill.fore_color.rgb = VERDE; par_num.line.fill.background(); par_num.shadow.inherit = False
    ptf = par_num.text_frame; ptf.word_wrap = False
    p(ptf, str(i + 1), 12, True, RGBColor(0xFF, 0xFF, 0xFF), first=True, align=PP_ALIGN.CENTER)
    tf = caja(s, 1.25, y - 0.04, 6.85, 0.85)
    par = tf.paragraphs[0]
    r = par.add_run(); r.text = t + '  '
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = VERDE_OSCURO; r.font.name = 'Arial'
    r2 = par.add_run(); r2.text = d
    r2.font.size = Pt(10.5); r2.font.color.rgb = TEXTO; r2.font.name = 'Arial'
    y += 1.0
marco(s, 8.35, 1.25, 2.25, 4.9, 'login.png',
      'Login con teclado visible (codigo de vendedor escrito)')
marco(s, 10.75, 1.25, 2.25, 4.9, 'dashboard.png',
      'Dashboard Hoy con clientes prioritarios')
tf = caja(s, 8.35, 6.35, 4.65, 0.8)
p(tf, 'La app recuerda tu sesión. Cierra sesión desde el menú de perfil si compartes el dispositivo.', 9.5, False, GRIS, first=True)

# ===== 5. FICHA COMERCIAL =====
s = nueva(); seccion(s, 'FICHA COMERCIAL Y EXPORT XLSX', 5)
tf = caja(s, 0.6, 1.0, 7.6, 0.75)
p(tf, 'La ficha concentra el contexto del cliente. El botón XLSX genera un libro de 5 hojas con fórmulas vivas (recalculan en Excel/LibreOffice):', 10.5, first=True)
tarjeta(s, 0.6, 1.9, 7.6, 1.5, 'Resumen · Ficha SKU', [
    'KPIs vivos: ventas, NCs, saldo neto, devuelto, SKUs · Evolución mensual con variación',
    'Todos los SKUs con moda 2026/2025 (xN) y fila TOTAL'])
tarjeta(s, 0.6, 3.55, 7.6, 1.5, 'Negociación · Comparativa', [
    'Anomalías: precio ≤15% bajo el promedio limpio · Análisis anual: moda, prom pond., min-max, Var %',
    'Pareto ~80% del saldo + ranking de precios por SKU'])
tf = caja(s, 0.6, 5.25, 7.6, 1.9)
p(tf, 'Cómo leer la comparativa', 12, True, VERDE_OSCURO, first=True, space=2)
p(tf, 'De quien paga más a quien paga menos: "Pos" es tu posición (RANK), "Δ vs mejor" cuánto pagas bajo respecto al mejor precio.', 10, space=3)
p(tf, 'Ejemplo (SKU 77264): MUNDO ESCOLAR S/ 50.20 (pos 1) → PRIMAVERA S/ 48.18 (pos 2, Δ −4.0%), precio congelado 2025-2026. Argumento: "te vengo dando el mejor precio de mi cartera".', 10, space=3)
p(tf, 'Anomalías: "SILICONA a S/ 0.66 vs tu promedio S/ 2.80 (−77%)": para reajuste o detectar error de facturación.', 10)
marco(s, 8.5, 1.2, 2.3, 4.85, 'ficha-precios.png',
      'Ficha PRIMAVERA con un SKU expandido (precio por ano)')
marco(s, 10.95, 1.2, 2.3, 4.85, 'ficha-comparativa.png',
      'Comparativa del SKU 77264, tu fila en negrita')
tf = caja(s, 8.5, 6.25, 4.75, 0.9)
p(tf, 'Opcional: captura del XLSX en Excel (hoja Negociación) como xlsx-negociacion.png.', 9.5, False, GRIS, first=True)

# ===== 6. INSTALACION PWA =====
s = nueva(); seccion(s, 'INSTALACIÓN PWA', 6)
tarjeta(s, 0.6, 1.15, 5.95, 1.9, 'Android (Chrome)', [
    'Abrir la URL de la app en Chrome',
    'Menú ⋮ → "Instalar aplicación"',
    'App standalone con icono propio + cache offline'])
tarjeta(s, 6.75, 1.15, 5.95, 1.9, 'iPhone / iPad (Safari)', [
    'Abrir la URL de la app en Safari',
    'Compartir → "Agregar a pantalla de inicio"',
    'Abre sin barras del navegador (standalone)'])
tf = caja(s, 0.6, 3.35, 12.1, 1.2)
p(tf, 'Nota de campo', 12.5, True, VERDE_OSCURO, first=True, space=3)
p(tf, 'Sin señal la app sirve cache (badges "Datos offline" u "obsoleto"). El stock se refresca ~15 min; el botón ↻ del Radar fuerza la actualización. Pull-to-refresh en Hoy, Radar y Clientes.', 11)
tf = caja(s, 0.6, 4.85, 12.1, 2.2)
p(tf, '¿Necesitas ayuda?', 12.5, True, VERDE_OSCURO, first=True, space=3)
p(tf, 'Si la app no carga o los datos no se actualizan: cierra y abre la app, revisa tu señal y vuelve a intentarlo.', 11, space=2)
p(tf, 'Si el problema persiste, reporta al administrador tu código de vendedor y el cliente que intentabas abrir.', 11)

# ===== 7. CONVENCIONES =====
s = nueva(); seccion(s, 'CONVENCIONES Y BUENAS PRÁCTICAS', 7)
conv = [
    ('Montos sin IGV', 'Los precios son de ERP (sin IGV). Toca un precio en el buscador para abrir la calculadora +18%.'),
    ('Códigos de display', 'Sin prefijo 01 (01186 → 186). Los SKUs nunca se recortan: 011019 ≠ 11019.'),
    ('Formato', 'es-PE, S/ con Intl.NumberFormat; cifras tabulares en tablas.'),
    ('Cache de 4 capas', 'memoria (15 min) → sessionStorage → localStorage → red. Cambiar de vendedor invalida el cache.'),
    ('Táctil', 'Botones ≥ 44px, inputs de 16px (anti-zoom iOS), filas completas tocables.'),
    ('Datos', 'Paginación offset con desempate folio_unico para no duplicar filas.'),
]
y = 1.15
for i, (t, d) in enumerate(conv):
    x = 0.6 if i % 2 == 0 else 6.75
    if i % 2 == 0 and i > 0: y += 1.85
    tarjeta(s, x, y, 5.95, 1.65, t, [d])

# ===== 8. LIMITES =====
s = nueva(); seccion(s, 'LÍMITES CONOCIDOS', 8)
lims = [
    ('Confidencialidad', 'Frenos de cliente (guardia de cartera 365d, bloqueo de login 3 min). Garantía real con Auth + RLS en Fase 4.'),
    ('Comparativa por SKU', 'Tope de 10k filas por SKU; SKUs de movimiento masivo pueden quedar parciales.'),
    ('Análisis anual', 'Solo SKUs con 3+ ventas en 4 años; baja rotación muestra último precio y modas.'),
    ('Historial grande', '+1000 filas caen a rebanadas mensuales; un mes fallido avisa "datos parciales".'),
    ('Netos B/C', 'Rangos muy largos pueden exceder el timeout: la app sugiere acortar el rango.'),
    ('XLSX', 'Las fórmulas vivas requieren Excel/LibreOffice; visores simples muestran valores cacheados (correctos).'),
]
y = 1.15
for i, (t, d) in enumerate(lims):
    x = 0.6 if i % 2 == 0 else 6.75
    if i % 2 == 0 and i > 0: y += 1.85
    barra = s.shapes.add_shape(1, Inches(x), Inches(y + 0.05), Inches(0.07), Inches(1.55))
    barra.fill.solid(); barra.fill.fore_color.rgb = AMBAR; barra.line.fill.background(); barra.shadow.inherit = False
    tf = caja(s, x + 0.25, y, 5.7, 1.75)
    p(tf, t, 11.5, True, AMBAR, first=True, space=3)
    p(tf, d, 10, False, TEXTO)

# ===== 9. ANEXO TECNICO =====
s = nueva(); seccion(s, 'ANEXO: DOCUMENTACIÓN TÉCNICA', 9)
tf = caja(s, 0.6, 1.0, 12.1, 0.5)
p(tf, 'Para el equipo de desarrollo. El deck de campo termina en el slide anterior.', 11, False, GRIS, first=True)
tarjeta(s, 0.6, 1.7, 5.95, 2.6, 'Desarrollo y deploy', [
    'npm install · npm run dev (vite) · npm run build (estático a build/)',
    'git push deploya a GitHub Pages vía GitHub Actions',
    'Base path: /g360-ventas-pulse',
    'OpenSpec: especificaciones y cambios en openspec/'])
tarjeta(s, 6.75, 1.7, 5.95, 2.6, 'Stack y variables', [
    'SvelteKit + adapter-static + Tailwind + ExcelJS (export) + IndexedDB (ruta) + Workbox (PWA)',
    'VITE_SUPABASE_URL · VITE_SUPABASE_ANON_KEY (solo lectura)',
    'VITE_STOCK_API_URL · VITE_STOCK_API_KEY'])

# ===== 10. RESUMEN =====
s = nueva(); seccion(s, 'RESUMEN', 10)
tf = caja(s, 0.6, 1.05, 12.1, 1.1)
p(tf, 'Ventas Pulse convierte el día del vendedor en un ciclo cerrado: prioridades claras al abrir, argumentos de precio con datos en la visita, y un export XLSX auditable para el seguimiento.', 13, False, TEXTO, first=True)
checks = [
    'Hoy: agenda priorizada con valor estimado y avance "Cómo voy"',
    'Radar + ruta del día persistente con stock en vivo',
    'Netos jerárquicos A/B/C con variación anual',
    'Ficha con anomalías de precio y comparativa Pareto',
    'Export XLSX de 5 hojas con fórmulas vivas y estilo CIPSA',
    'Búsqueda global con voz (SKU dígito a dígito o por nombre)',
    'Offline-first con cache de 4 capas y PWA instalable',
]
y = 2.35
for c in checks:
    tf = caja(s, 1.0, y, 11.5, 0.45)
    par = tf.paragraphs[0]
    r = par.add_run(); r.text = '✓  '
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = VERDE; r.font.name = 'Arial'
    r2 = par.add_run(); r2.text = c
    r2.font.size = Pt(11.5); r2.font.color.rgb = TEXTO; r2.font.name = 'Arial'
    y += 0.56
tf = caja(s, 0.6, 6.6, 12.1, 0.4)
p(tf, 'Documentación completa: README.md · Especificaciones: openspec/ · powered by G360', 10, False, GRIS, first=True, align=PP_ALIGN.CENTER)

prs.save('manual-ventas-pulse.pptx')
print(f'OK manual-ventas-pulse.pptx — {TOTAL} slides 16:9 {W}x{H}in')
